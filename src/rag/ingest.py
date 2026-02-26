from __future__ import annotations

import io
import json
import re
from typing import Any, Dict, List, Tuple


def _parse_text_payload(payload: bytes) -> List[Dict[str, Any]]:
    text = payload.decode("utf-8", errors="ignore").strip()
    if not text:
        return []
    return [{"text": text, "metadata": {"section": "full_text"}}]


def _parse_pdf_payload(payload: bytes) -> List[Dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover - import availability
        raise RuntimeError("缺少 pypdf 依赖，请先安装 requirements.txt") from exc

    reader = PdfReader(io.BytesIO(payload))
    segments: List[Dict[str, Any]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        segments.append({"text": text, "metadata": {"page": idx}})
    return segments


def _parse_ppt_payload(payload: bytes) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - import availability
        raise RuntimeError("缺少 python-pptx 依赖，请先安装 requirements.txt") from exc

    prs = Presentation(io.BytesIO(payload))
    segments: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                val = str(shape.text or "").strip()
                if val:
                    lines.append(val)
        if lines:
            segments.append({"text": "\n".join(lines), "metadata": {"slide": idx}})
    return segments


def _flatten_rich_text(value: Any) -> str:
    if isinstance(value, str):
        return value
    if isinstance(value, list):
        out: List[str] = []
        for item in value:
            if isinstance(item, str):
                out.append(item)
            elif isinstance(item, dict):
                part = item.get("text")
                if isinstance(part, str):
                    out.append(part)
        return "".join(out)
    return ""


def _parse_telegram_json(payload: bytes) -> List[Dict[str, Any]]:
    data = json.loads(payload.decode("utf-8", errors="ignore"))
    messages = data.get("messages", []) if isinstance(data, dict) else []
    if not isinstance(messages, list):
        return []

    segments: List[Dict[str, Any]] = []
    for idx, msg in enumerate(messages, start=1):
        if not isinstance(msg, dict):
            continue
        text = _flatten_rich_text(msg.get("text")).strip()
        if not text:
            continue
        author = str(msg.get("from") or msg.get("actor") or "unknown")
        dt = str(msg.get("date") or msg.get("date_unixtime") or "")
        segments.append(
            {
                "text": f"[{dt}] {author}: {text}",
                "metadata": {
                    "message_index": idx,
                    "author": author,
                    "date": dt,
                },
            }
        )
    return segments


def _parse_generic_chat_json(payload: bytes) -> List[Dict[str, Any]]:
    data = json.loads(payload.decode("utf-8", errors="ignore"))
    items = data
    if isinstance(data, dict):
        items = data.get("messages", data.get("items", []))
    if not isinstance(items, list):
        return []

    segments: List[Dict[str, Any]] = []
    for idx, msg in enumerate(items, start=1):
        if not isinstance(msg, dict):
            continue
        author = str(msg.get("from") or msg.get("sender") or msg.get("role") or "unknown")
        text = str(msg.get("text") or msg.get("content") or msg.get("message") or "").strip()
        dt = str(msg.get("date") or msg.get("timestamp") or "")
        if not text:
            continue
        segments.append(
            {
                "text": f"[{dt}] {author}: {text}",
                "metadata": {
                    "message_index": idx,
                    "author": author,
                    "date": dt,
                },
            }
        )
    return segments


def _parse_whatsapp_text(payload: bytes) -> List[Dict[str, Any]]:
    text = payload.decode("utf-8", errors="ignore")
    lines = [x.strip() for x in text.splitlines() if x.strip()]
    segments: List[Dict[str, Any]] = []

    # Common WhatsApp export patterns:
    # [12/01/24, 10:30:12 AM] Alice: hello
    # 12/01/24, 10:30 - Alice: hello
    patterns = [
        re.compile(r"^\[(?P<date>.+?)\]\s(?P<author>[^:]+):\s(?P<msg>.+)$"),
        re.compile(r"^(?P<date>\d{1,2}/\d{1,2}/\d{2,4},?\s.+?)\s-\s(?P<author>[^:]+):\s(?P<msg>.+)$"),
    ]

    for idx, line in enumerate(lines, start=1):
        matched = None
        for pat in patterns:
            m = pat.match(line)
            if m:
                matched = m
                break
        if not matched:
            continue
        dt = matched.group("date").strip()
        author = matched.group("author").strip()
        msg = matched.group("msg").strip()
        if not msg:
            continue
        segments.append(
            {
                "text": f"[{dt}] {author}: {msg}",
                "metadata": {"message_index": idx, "author": author, "date": dt, "platform": "whatsapp"},
            }
        )
    return segments


def _parse_whatsapp_json(payload: bytes) -> List[Dict[str, Any]]:
    data = json.loads(payload.decode("utf-8", errors="ignore"))
    items = data.get("messages", data) if isinstance(data, dict) else data
    if not isinstance(items, list):
        return []
    segments: List[Dict[str, Any]] = []
    for idx, msg in enumerate(items, start=1):
        if not isinstance(msg, dict):
            continue
        author = str(msg.get("sender") or msg.get("from") or msg.get("author") or "unknown")
        content = str(msg.get("text") or msg.get("message") or msg.get("content") or "").strip()
        dt = str(msg.get("timestamp") or msg.get("date") or msg.get("time") or "")
        if not content:
            continue
        segments.append(
            {
                "text": f"[{dt}] {author}: {content}",
                "metadata": {"message_index": idx, "author": author, "date": dt, "platform": "whatsapp"},
            }
        )
    return segments


def _parse_discord_json(payload: bytes) -> List[Dict[str, Any]]:
    data = json.loads(payload.decode("utf-8", errors="ignore"))
    items = data.get("messages", data.get("data", data)) if isinstance(data, dict) else data
    if not isinstance(items, list):
        return []
    segments: List[Dict[str, Any]] = []
    for idx, msg in enumerate(items, start=1):
        if not isinstance(msg, dict):
            continue
        author_obj = msg.get("author") if isinstance(msg.get("author"), dict) else {}
        author = str(
            author_obj.get("name")
            or author_obj.get("username")
            or msg.get("authorName")
            or msg.get("sender")
            or "unknown"
        )
        content = str(msg.get("content") or msg.get("text") or "").strip()
        dt = str(msg.get("timestamp") or msg.get("createdAt") or msg.get("date") or "")
        if not content:
            continue
        segments.append(
            {
                "text": f"[{dt}] {author}: {content}",
                "metadata": {"message_index": idx, "author": author, "date": dt, "platform": "discord"},
            }
        )
    return segments


def _parse_slack_json(payload: bytes) -> List[Dict[str, Any]]:
    data = json.loads(payload.decode("utf-8", errors="ignore"))
    items = data.get("messages", data) if isinstance(data, dict) else data
    if not isinstance(items, list):
        return []
    segments: List[Dict[str, Any]] = []
    for idx, msg in enumerate(items, start=1):
        if not isinstance(msg, dict):
            continue
        author = str(msg.get("user_profile", {}).get("real_name") if isinstance(msg.get("user_profile"), dict) else "") or str(
            msg.get("username") or msg.get("user") or msg.get("bot_id") or "unknown"
        )
        content = str(msg.get("text") or msg.get("content") or "").strip()
        dt = str(msg.get("ts") or msg.get("timestamp") or msg.get("date") or "")
        if not content:
            continue
        segments.append(
            {
                "text": f"[{dt}] {author}: {content}",
                "metadata": {"message_index": idx, "author": author, "date": dt, "platform": "slack"},
            }
        )
    return segments


def parse_uploaded_file(file_name: str, payload: bytes, import_type: str) -> Tuple[str, List[Dict[str, Any]]]:
    lower = str(file_name or "").lower()
    mode = str(import_type or "document").strip().lower()

    if mode in {"telegram", "telegram-chat", "telegram_chat"}:
        if lower.endswith(".json"):
            return "telegram", _parse_telegram_json(payload)
        return "telegram", _parse_text_payload(payload)

    if mode in {"chat", "third-party-chat", "third_party_chat"}:
        if lower.endswith(".json"):
            return "third_party_chat", _parse_generic_chat_json(payload)
        return "third_party_chat", _parse_text_payload(payload)

    if mode in {"whatsapp", "whatsapp-chat", "whatsapp_chat"}:
        if lower.endswith(".json"):
            return "whatsapp", _parse_whatsapp_json(payload)
        return "whatsapp", _parse_whatsapp_text(payload)

    if mode in {"discord", "discord-chat", "discord_chat"}:
        if lower.endswith(".json"):
            return "discord", _parse_discord_json(payload)
        return "discord", _parse_text_payload(payload)

    if mode in {"slack", "slack-chat", "slack_chat"}:
        if lower.endswith(".json"):
            return "slack", _parse_slack_json(payload)
        return "slack", _parse_text_payload(payload)

    if lower.endswith(".pdf"):
        return "document", _parse_pdf_payload(payload)
    if lower.endswith(".pptx") or lower.endswith(".ppt"):
        return "document", _parse_ppt_payload(payload)
    if lower.endswith(".json"):
        return "document", _parse_generic_chat_json(payload)
    return "document", _parse_text_payload(payload)
